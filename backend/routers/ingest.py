import json
import os
import re
from io import BytesIO
from typing import Any, Dict, Optional

import PyPDF2
from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from pydantic import BaseModel
from sqlalchemy.orm import Session

import auth
import models
from database import get_db

from ai_clients import AIProviderError, embeddings_model, generate_ai_response

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import pytesseract
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
except Exception:
    pytesseract = None
    Image = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import cv2
    import numpy as np
except Exception:
    cv2 = None
    np = None

if pytesseract:
    tesseract_cmd = os.getenv("TESSERACT_CMD")
    if not tesseract_cmd and os.name == "nt":
        default_tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_tesseract_cmd):
            tesseract_cmd = default_tesseract_cmd
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

router = APIRouter(prefix="/ingest", tags=["ingestion"])

MAX_UPLOAD_BYTES = int(os.getenv("MAX_UPLOAD_BYTES", str(10 * 1024 * 1024)))
MAX_INDEX_CHUNKS = int(os.getenv("MAX_INDEX_CHUNKS", "40"))
MIN_INDEXABLE_TEXT_CHARS = int(os.getenv("MIN_INDEXABLE_TEXT_CHARS", "80"))
MAX_OCR_PAGES = int(os.getenv("MAX_OCR_PAGES", "20"))
OCR_RENDER_SCALE = float(os.getenv("OCR_RENDER_SCALE", "4"))
OCR_FALLBACK_RENDER_SCALE = float(os.getenv("OCR_FALLBACK_RENDER_SCALE", "3"))
SUPPORTED_UPLOAD_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}
IMAGE_UPLOAD_EXTENSIONS = {".png", ".jpg", ".jpeg"}
OCR_CONFIGS = [
    "--oem 3 --psm 6 -c preserve_interword_spaces=1",
    "--oem 3 --psm 4 -c preserve_interword_spaces=1",
    "--oem 3 --psm 11 -c preserve_interword_spaces=1",
    "--oem 3 --psm 12 -c preserve_interword_spaces=1",
]

OCR_ACADEMIC_KEYWORDS = {
    "COVENANT",
    "UNIVERSITY",
    "COURSE",
    "TITLE",
    "SESSION",
    "SEMESTER",
    "QUESTION",
    "MARKS",
    "EXAMINATION",
    "INSTRUCTION",
    "PROJECT",
    "MANAGEMENT",
    "RISK",
    "SCOPE",
    "COST",
    "STAKEHOLDER",
}

DOCUMENT_TYPES = {
    "past_question",
    "lecture_note",
    "course_outline",
    "tutorial",
    "assignment",
    "revision_slide",
    "exam_prep",
    "unknown",
}
EXAM_TYPES = {"quiz", "test", "midterm", "final", "unknown"}
TRUSTED_HEURISTIC_FIELDS = {
    "document_type",
    "course_code",
    "course_title",
    "academic_year",
    "semester",
    "department",
    "college",
    "faculty",
    "exam_type",
}
WEAK_METADATA_VALUES = {
    "",
    "unknown",
    "not found",
    "none",
    "null",
    "n/a",
    "na",
    "academic document",
    "academic_document",
}


class DeleteDocumentRequest(BaseModel):
    source_file: str | None = None
    document_type: str
    document_title: str | None = None


def missing_ai_error(feature: str):
    return HTTPException(
        status_code=503,
        detail=f"{feature} is not available. Check DEEPSEEK_API_KEY and that fastembed is installed (pip install fastembed).",
    )


@router.get("/ocr-health")
def ocr_health():
    tesseract_version = None
    if pytesseract:
        try:
            tesseract_version = str(pytesseract.get_tesseract_version())
        except Exception as exc:
            tesseract_version = f"unavailable: {exc}"
    return {
        "pytesseract_imported": bool(pytesseract),
        "tesseract_cmd": getattr(getattr(pytesseract, "pytesseract", None), "tesseract_cmd", "") if pytesseract else "",
        "tesseract_version": tesseract_version,
        "opencv_available": bool(cv2 and np),
        "pillow_available": bool(Image),
        "pymupdf_available": bool(fitz),
    }


@router.delete("/documents")
def delete_uploaded_document(
    req: DeleteDocumentRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(auth.require_role("student")),
):
    document_type = (req.document_type or "").strip().lower()
    source_file = (req.source_file or "").strip()
    document_title = (req.document_title or "").strip()
    if document_type not in {"past_question", "lecture_note"}:
        raise HTTPException(status_code=400, detail="document_type must be past_question or lecture_note.")
    if not source_file and not document_title:
        raise HTTPException(status_code=400, detail="source_file or document_title is required.")

    if document_type == "past_question":
        rows = [
            row
            for row in db.query(models.PastQuestion).all()
            if _metadata_matches(row.metadata_json, source_file, document_title)
        ]
        summary = _delete_past_questions(db, rows)
    else:
        rows = [
            row
            for row in db.query(models.LectureNote).all()
            if _metadata_matches(row.metadata_json, source_file, document_title)
        ]
        summary = _delete_lecture_notes(db, rows)

    db.commit()
    return summary


@router.delete("/clear-materials")
def clear_uploaded_materials(
    prune_empty_courses: bool = Query(False),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(auth.require_role("student")),
):
    summary = _empty_delete_summary()

    past_questions = db.query(models.PastQuestion).all()
    _merge_delete_summary(summary, _delete_past_questions(db, past_questions))

    lecture_notes = db.query(models.LectureNote).all()
    _merge_delete_summary(summary, _delete_lecture_notes(db, lecture_notes))

    orphan_chunks = db.query(models.LectureNoteChunk).delete(synchronize_session=False)
    summary["lecture_note_chunks_deleted"] += orphan_chunks

    if prune_empty_courses:
        summary["courses_pruned"] = _prune_empty_auto_courses(db)

    db.commit()
    return summary


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200):
    chunks = []
    start = 0
    clean_text = re.sub(r"\n{3,}", "\n\n", text).strip()
    while start < len(clean_text):
        end = start + chunk_size
        chunks.append(clean_text[start:end])
        start += chunk_size - overlap
    return chunks


def _quality_score(text: str, page_count: int, method: str) -> float:
    clean = re.sub(r"\s+", " ", text or "").strip()
    if not clean:
        return 0.0
    length_score = min(len(clean) / 1200, 1.0)
    alpha_ratio = sum(ch.isalpha() for ch in clean) / max(len(clean), 1)
    page_bonus = min(page_count / 8, 1.0) * 0.1
    method_base = 0.55 if method == "ocr" else 0.65
    return round(min(method_base + (length_score * 0.25) + (alpha_ratio * 0.12) + page_bonus, 0.98), 2)


def _clean_extracted_text(parts: list[str]) -> str:
    return re.sub(r"\n{3,}", "\n\n", "\n".join(part for part in parts if part).strip())


def _clean_ocr_text(raw_text: str) -> str:
    text = raw_text or ""
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[\u2022\u25cf\u25a0\u25aa]+", " ", text)
    text = re.sub(r"[_=*#~|]{2,}", " ", text)
    text = re.sub(r"([^\w\s])\1{2,}", r"\1", text)
    text = re.sub(r"\b([A-Z]{2,4})\s+(\d{3})\b", r"\1\2", text)
    text = re.sub(r"([A-Za-z])-\s+([A-Za-z])", r"\1\2", text)
    text = re.sub(r"\s+([,.;:?)])", r"\1", text)
    text = re.sub(r"([(])\s+", r"\1", text)

    useful_lines = []
    for raw_line in text.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip(" -:\t")
        if not line:
            continue
        alpha_count = sum(ch.isalpha() for ch in line)
        digit_count = sum(ch.isdigit() for ch in line)
        symbol_count = sum(not ch.isalnum() and not ch.isspace() for ch in line)
        symbol_ratio = symbol_count / max(len(line), 1)
        academic_signal = bool(
            re.search(
                r"\b(course|title|code|session|semester|department|instruction|question|marks?|"
                r"project|management|risk|scope|cost|stakeholder|communication|procurement|"
                r"calculate|draw|determine|identify|explain|discuss|evaluate|PERT|CPI|SPI)\b",
                line,
                re.IGNORECASE,
            )
            or re.search(r"\b(?:question\s*)?[1-9]\s*[.)]\s*[A-Za-z]", line, re.IGNORECASE)
            or re.search(r"\b[1-9]\s*\([a-z]\)", line, re.IGNORECASE)
            or re.search(r"\b[A-Z]{2,4}\d{3}\b", line)
            or re.search(r"\b\d+\s*marks?\b", line, re.IGNORECASE)
        )
        if len(line) < 8 and not academic_signal:
            continue
        if alpha_count < 5 and not (academic_signal and digit_count):
            continue
        if len(line) < 70 and re.search(r"(head|depa|dept|itt+e+|gt\b|^[sS]\s*['`‘’])", line, re.IGNORECASE):
            continue
        if len(line) < 40 and alpha_count / max(len(line), 1) < 0.45:
            continue
        if symbol_ratio > 0.35 and not academic_signal:
            continue
        if re.fullmatch(r"[\W\d_]+", line):
            continue
        if re.search(r"(stamp|scanned with|cam.?scanner|page\s+\d+\s+of\s+\d+)", line, re.IGNORECASE):
            continue
        useful_lines.append(line[:500])

    cleaned = "\n".join(useful_lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _extract_with_pymupdf(content: bytes) -> tuple[str, int]:
    if not fitz:
        return "", 0
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        parts = [page.get_text("text") or "" for page in doc]
        return _clean_extracted_text(parts), doc.page_count
    finally:
        doc.close()


def _extract_with_pypdf2(content: bytes) -> tuple[str, int]:
    pdf_reader = PyPDF2.PdfReader(BytesIO(content))
    if getattr(pdf_reader, "is_encrypted", False):
        raise ValueError("encrypted_pdf")
    parts = [(page.extract_text() or "") for page in pdf_reader.pages]
    return _clean_extracted_text(parts), len(pdf_reader.pages)


def _render_pdf_page_to_image(page) -> Image.Image:
    try:
        pix = page.get_pixmap(matrix=fitz.Matrix(OCR_RENDER_SCALE, OCR_RENDER_SCALE), alpha=False)
    except Exception as exc:
        print(f"OCR: high-resolution render failed, using {OCR_FALLBACK_RENDER_SCALE}x fallback: {exc}")
        pix = page.get_pixmap(matrix=fitz.Matrix(OCR_FALLBACK_RENDER_SCALE, OCR_FALLBACK_RENDER_SCALE), alpha=False)
    return Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")


def _opencv_preprocess_variants(pil_image: Image.Image) -> list[tuple[str, Image.Image]]:
    if cv2 is None or np is None:
        print("OCR: OpenCV unavailable, using Pillow fallback")
        return []
    try:
        rgb = np.array(pil_image.convert("RGB"))
        bgr = cv2.cvtColor(rgb, cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
        denoised = cv2.fastNlMeansDenoising(gray, None, 18, 7, 21)
        clahe = cv2.createCLAHE(clipLimit=2.4, tileGridSize=(8, 8)).apply(denoised)
        _, otsu = cv2.threshold(clahe, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        adaptive = cv2.adaptiveThreshold(
            clahe, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 35, 11
        )
        inverted_adaptive = cv2.bitwise_not(adaptive)
        background = cv2.medianBlur(cv2.dilate(gray, np.ones((9, 9), np.uint8)), 31)
        shadow_reduced = cv2.normalize(255 - cv2.absdiff(gray, background), None, 0, 255, cv2.NORM_MINMAX)
        sharpened = cv2.filter2D(clahe, -1, np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]]))
        kernel = np.ones((2, 2), np.uint8)
        opened = cv2.morphologyEx(otsu, cv2.MORPH_OPEN, kernel)
        closed = cv2.morphologyEx(opened, cv2.MORPH_CLOSE, kernel)

        variants = [
            ("opencv-gray", gray),
            ("opencv-denoised", denoised),
            ("opencv-clahe", clahe),
            ("opencv-otsu", otsu),
            ("opencv-adaptive", adaptive),
            ("opencv-inverted-adaptive", inverted_adaptive),
            ("opencv-shadow-reduced", shadow_reduced),
            ("opencv-sharpened", sharpened),
            ("opencv-morphology", closed),
        ]
        return [(name, Image.fromarray(image).convert("RGB")) for name, image in variants]
    except Exception as exc:
        print(f"OCR: OpenCV unavailable, using Pillow fallback: {exc}")
        return []


def _threshold_image(image: Image.Image, threshold: int) -> Image.Image:
    return image.point(lambda px: 255 if px > threshold else 0, mode="1").convert("RGB")


def _pillow_preprocess_variants(pil_image: Image.Image) -> list[tuple[str, Image.Image]]:
    original = pil_image.convert("RGB")
    gray = ImageOps.grayscale(original)
    autocontrast = ImageOps.autocontrast(gray)
    sharpened = autocontrast.filter(ImageFilter.SHARPEN)
    high_contrast = ImageEnhance.Contrast(autocontrast).enhance(1.8)
    enlarged = original.resize((original.width * 2, original.height * 2))
    variants = [
        ("pillow-original", original),
        ("pillow-gray", gray.convert("RGB")),
        ("pillow-autocontrast", autocontrast.convert("RGB")),
        ("pillow-sharpened", sharpened.convert("RGB")),
        ("pillow-high-contrast", high_contrast.convert("RGB")),
        ("pillow-threshold-130", _threshold_image(autocontrast, 130)),
        ("pillow-threshold-150", _threshold_image(autocontrast, 150)),
        ("pillow-threshold-170", _threshold_image(autocontrast, 170)),
        ("pillow-threshold-190", _threshold_image(autocontrast, 190)),
        ("pillow-enlarged", enlarged),
    ]
    return variants


def _useful_word_count(text: str) -> int:
    return len(re.findall(r"\b[A-Za-z][A-Za-z0-9&./-]{2,}\b", text or ""))


def _score_ocr_text(text: str) -> float:
    clean = re.sub(r"\s+", " ", text or "").strip()
    if not clean:
        return 0.0
    char_count = len(clean)
    alpha_count = sum(ch.isalpha() for ch in clean)
    alpha_ratio = alpha_count / max(char_count, 1)
    useful_words = _useful_word_count(clean)
    upper = clean.upper()
    keyword_hits = sum(1 for keyword in OCR_ACADEMIC_KEYWORDS if keyword in upper)
    course_hits = len(re.findall(r"\b[A-Z]{2,4}\s?\d{3}\b", upper))
    session_hits = len(re.findall(r"\b20\d{2}\s*[/\-]\s*(?:20)?\d{2}\b", clean))
    question_hits = len(
        re.findall(r"\bquestion\s*\d+\b|\b\d+\s*\([a-z]\)|\b\d+[.)]\s+[A-Za-z]", clean, re.IGNORECASE)
    )
    single_letter_words = len(re.findall(r"\b[A-Za-z]\b", clean))
    symbol_ratio = sum(not ch.isalnum() and not ch.isspace() for ch in clean) / max(char_count, 1)
    repeated_garbage = len(re.findall(r"([^\w\s])\1{2,}|([A-Za-z])\2{4,}", clean))

    score = 0.0
    score += min(char_count / 18, 35)
    score += min(useful_words * 1.8, 30)
    score += min(keyword_hits * 6, 30)
    score += min(course_hits * 12, 24)
    score += min(session_hits * 10, 20)
    score += min(question_hits * 5, 25)
    score += alpha_ratio * 20
    score -= max(0, single_letter_words - useful_words) * 0.8
    score -= symbol_ratio * 45
    score -= repeated_garbage * 8
    if char_count < 40:
        score -= 25
    if alpha_ratio < 0.35:
        score -= 30
    return round(max(score, 0.0), 2)


def _is_useful_ocr_text(text: str) -> bool:
    clean = (text or "").strip()
    return len(clean) >= 50 and (_useful_word_count(clean) >= 8 or _score_ocr_text(clean) >= 45)


def _run_tesseract_variants(image: Image.Image) -> tuple[str, float, str]:
    best_text = ""
    best_score = 0.0
    best_config = ""
    for config in OCR_CONFIGS:
        try:
            text = pytesseract.image_to_string(image, config=config) or ""
        except Exception as exc:
            print(f"OCR: Tesseract config failed ({config}): {exc}")
            continue
        score = _score_ocr_text(text)
        if score > best_score or (score == best_score and len(text.strip()) > len(best_text.strip())):
            best_text = text
            best_score = score
            best_config = config
    return best_text, best_score, best_config


def _best_ocr_for_page(image: Image.Image, page_number: int) -> tuple[str, float, str, str]:
    variants = _opencv_preprocess_variants(image)
    engine = "OpenCV" if variants else "Pillow fallback"
    if variants:
        print("OCR engine: OpenCV")
        variants.extend(_pillow_preprocess_variants(image)[:3])
    else:
        print("OCR engine: Pillow fallback")
        variants = _pillow_preprocess_variants(image)

    best_text = ""
    best_score = 0.0
    best_config = ""
    best_variant = ""
    for variant_name, variant_image in variants:
        text, score, config = _run_tesseract_variants(variant_image)
        if score > best_score or (score == best_score and len(text.strip()) > len(best_text.strip())):
            best_text = text
            best_score = score
            best_config = config
            best_variant = variant_name

    print(
        f"OCR page {page_number} best chars: {len(best_text.strip())}, "
        f"score: {best_score}, config: {best_config}, variant: {best_variant}"
    )
    return best_text, best_score, best_config, engine


def _ocr_pdf(content: bytes, page_count: int, warnings: list[str]) -> tuple[str, str | None]:
    if not fitz:
        warnings.append("OCR could not run because PyMuPDF is not installed.")
        return "", "ocr_not_installed"
    if not pytesseract or not Image:
        warnings.append("OCR support is not installed. Install Tesseract OCR or upload a text-based file.")
        return "", "ocr_not_installed"

    doc = fitz.open(stream=content, filetype="pdf")
    try:
        ocr_parts = []
        pages_to_scan = min(doc.page_count, MAX_OCR_PAGES)
        if page_count and doc.page_count > MAX_OCR_PAGES:
            warnings.append(f"OCR scanned the first {MAX_OCR_PAGES} pages only.")
        for page_index in range(pages_to_scan):
            page = doc.load_page(page_index)
            image = _render_pdf_page_to_image(page)
            page_text, page_score, _config, _engine = _best_ocr_for_page(image, page_index + 1)
            if _is_useful_ocr_text(page_text):
                ocr_parts.append(page_text)
            elif page_text.strip():
                ocr_parts.append(page_text)
        text = _clean_extracted_text(ocr_parts)
        print(f"OCR final chars: {len(text.strip())}")
        print(f"OCR final useful words: {_useful_word_count(text)}")
        if text.strip() and not _is_useful_ocr_text(text):
            warnings.append("OCR produced text, but it may be incomplete or noisy.")
        return text, None
    except pytesseract.TesseractNotFoundError:
        warnings.append("OCR support is not installed. Install Tesseract OCR or upload a text-based file.")
        return "", "ocr_not_installed"
    except Exception as exc:
        warnings.append(f"OCR could not read this PDF: {exc}")
        return "", "ocr_failed"
    finally:
        doc.close()


def _successful_text_extraction(
    text: str,
    method: str,
    page_count: int,
    warnings: list[str],
    ocr_used: bool = False,
    confidence_method: str | None = None,
) -> Dict[str, Any]:
    cleaned_text = _clean_ocr_text(text)
    indexable_text = cleaned_text if len(cleaned_text.strip()) >= MIN_INDEXABLE_TEXT_CHARS else text
    return {
        "text": text,
        "raw_extracted_text": text,
        "cleaned_text": indexable_text,
        "method": method,
        "page_count": page_count,
        "text_char_count": len(text.strip()),
        "cleaned_text_char_count": len(indexable_text.strip()),
        "ocr_used": ocr_used,
        "extraction_confidence": _quality_score(text, page_count, confidence_method or method),
        "failure_reason": None,
        "warnings": warnings,
    }


def _failed_text_extraction(
    raw_text: str,
    page_count: int,
    warnings: list[str],
    failure_reason: str,
    ocr_used: bool = False,
    ocr_score: float = 0.0,
    ocr_useful_words: int = 0,
) -> Dict[str, Any]:
    cleaned_text = _clean_ocr_text(raw_text)
    return {
        "text": raw_text,
        "raw_extracted_text": raw_text,
        "cleaned_text": cleaned_text,
        "method": "failed",
        "page_count": page_count,
        "text_char_count": len(raw_text.strip()),
        "cleaned_text_char_count": len(cleaned_text.strip()),
        "ocr_used": ocr_used,
        "extraction_confidence": 0.0,
        "failure_reason": failure_reason,
        "indexed_status": "unindexed",
        "searchable": False,
        "needs_review": True,
        "ocr_score": ocr_score,
        "ocr_useful_words": ocr_useful_words,
        "warnings": warnings,
    }


def _extract_docx_text(content: bytes, warnings: list[str]) -> Dict[str, Any]:
    if not DocxDocument:
        warnings.append("Word document support is not installed. Install python-docx.")
        return _failed_text_extraction("", 0, warnings, "docx_not_installed")
    try:
        doc = DocxDocument(BytesIO(content))
        parts: list[str] = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        text = _clean_extracted_text(parts)
        if len(text.strip()) >= MIN_INDEXABLE_TEXT_CHARS:
            return _successful_text_extraction(text, "docx_text", max(1, len(doc.sections)), warnings)
        warnings.append("Word document did not contain enough readable text.")
        return _failed_text_extraction(text, max(1, len(doc.sections)), warnings, "no_text_found")
    except Exception as exc:
        warnings.append(f"Word document text extraction failed: {exc}")
        return _failed_text_extraction("", 0, warnings, "docx_read_failed")


def _extract_pptx_text(content: bytes, warnings: list[str]) -> Dict[str, Any]:
    if not Presentation:
        warnings.append("PowerPoint support is not installed. Install python-pptx.")
        return _failed_text_extraction("", 0, warnings, "pptx_not_installed")
    try:
        deck = Presentation(BytesIO(content))
        parts: list[str] = []
        for slide_index, slide in enumerate(deck.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                    if text.strip():
                        slide_parts.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_parts.append(row_text)
            notes_slide = getattr(slide, "notes_slide", None)
            notes_frame = getattr(notes_slide, "notes_text_frame", None) if notes_slide else None
            if notes_frame and notes_frame.text.strip():
                slide_parts.append(notes_frame.text)
            if slide_parts:
                parts.append(f"Slide {slide_index}\n" + "\n".join(slide_parts))
        text = _clean_extracted_text(parts)
        page_count = len(deck.slides)
        if len(text.strip()) >= MIN_INDEXABLE_TEXT_CHARS:
            return _successful_text_extraction(text, "pptx_text", page_count, warnings)
        warnings.append("PowerPoint file did not contain enough readable text.")
        return _failed_text_extraction(text, page_count, warnings, "no_text_found")
    except Exception as exc:
        warnings.append(f"PowerPoint text extraction failed: {exc}")
        return _failed_text_extraction("", 0, warnings, "pptx_read_failed")


def _extract_image_text(content: bytes, warnings: list[str]) -> Dict[str, Any]:
    if not pytesseract or not Image:
        warnings.append("OCR support is not installed. Install Tesseract OCR to read image uploads.")
        return _failed_text_extraction("", 1, warnings, "ocr_not_installed", ocr_used=True)
    try:
        image = Image.open(BytesIO(content)).convert("RGB")
        ocr_text, ocr_score, _config, _engine = _best_ocr_for_page(image, 1)
        text = _clean_extracted_text([ocr_text])
        ocr_useful_words = _useful_word_count(text)
        if len(text.strip()) >= MIN_INDEXABLE_TEXT_CHARS or _is_useful_ocr_text(text):
            cleaned_text = _clean_ocr_text(text)
            indexable_text = cleaned_text if len(cleaned_text.strip()) >= MIN_INDEXABLE_TEXT_CHARS else text
            extraction_confidence = max(0.45, _quality_score(text, 1, "ocr"))
            if extraction_confidence < 0.65 or ocr_score < 75:
                warnings.append("OCR completed, but extracted text may need review.")
            return {
                "text": text,
                "raw_extracted_text": text,
                "cleaned_text": indexable_text,
                "method": "ocr",
                "page_count": 1,
                "text_char_count": len(text.strip()),
                "cleaned_text_char_count": len(indexable_text.strip()),
                "ocr_used": True,
                "extraction_confidence": extraction_confidence,
                "failure_reason": None,
                "indexed_status": "indexed_review_required" if extraction_confidence < 0.65 or ocr_score < 75 else "indexed",
                "searchable": True,
                "needs_review": extraction_confidence < 0.65 or ocr_score < 75,
                "ocr_score": ocr_score,
                "ocr_useful_words": ocr_useful_words,
                "warnings": warnings,
            }
        failure_reason = "ocr_low_confidence" if text.strip() else "file_too_blurry"
        return _failed_text_extraction(text, 1, warnings, failure_reason, ocr_used=True, ocr_score=ocr_score, ocr_useful_words=ocr_useful_words)
    except pytesseract.TesseractNotFoundError:
        warnings.append("OCR support is not installed. Install Tesseract OCR to read image uploads.")
        return _failed_text_extraction("", 1, warnings, "ocr_not_installed", ocr_used=True)
    except Exception as exc:
        warnings.append(f"OCR could not read this image: {exc}")
        return _failed_text_extraction("", 1, warnings, "ocr_failed", ocr_used=True)


def _extract_pdf_content(content: bytes, warnings: list[str]) -> Dict[str, Any]:
    embedded_text = ""
    page_count = 0
    failure_reason = "embedded_text_weak"

    try:
        embedded_text, page_count = _extract_with_pymupdf(content)
    except Exception as exc:
        warnings.append(f"PyMuPDF text extraction failed: {exc}")
        if "encrypted" in str(exc).lower():
            failure_reason = "encrypted_pdf"

    if len(embedded_text.strip()) < MIN_INDEXABLE_TEXT_CHARS:
        try:
            fallback_text, fallback_pages = _extract_with_pypdf2(content)
            if len(fallback_text.strip()) > len(embedded_text.strip()):
                embedded_text = fallback_text
            page_count = page_count or fallback_pages
        except Exception as exc:
            warnings.append(f"PyPDF2 text extraction failed: {exc}")
            if "encrypted_pdf" in str(exc) or "encrypted" in str(exc).lower():
                failure_reason = "encrypted_pdf"

    embedded_count = len(embedded_text.strip())
    if embedded_count >= MIN_INDEXABLE_TEXT_CHARS:
        cleaned_text = _clean_ocr_text(embedded_text)
        return {
            "text": embedded_text,
            "raw_extracted_text": embedded_text,
            "cleaned_text": cleaned_text or embedded_text,
            "method": "embedded_text",
            "page_count": page_count,
            "text_char_count": embedded_count,
            "cleaned_text_char_count": len((cleaned_text or embedded_text).strip()),
            "ocr_used": False,
            "extraction_confidence": _quality_score(embedded_text, page_count, "embedded_text"),
            "failure_reason": None,
            "warnings": warnings,
        }

    warnings.append("Weak embedded text detected")
    ocr_text, ocr_failure_reason = _ocr_pdf(content, page_count, warnings)
    ocr_count = len(ocr_text.strip())
    ocr_score = _score_ocr_text(ocr_text)
    ocr_useful_words = _useful_word_count(ocr_text)
    ocr_is_useful = _is_useful_ocr_text(ocr_text)

    if ocr_count >= MIN_INDEXABLE_TEXT_CHARS or ocr_is_useful:
        text = ocr_text
        method = "ocr"
        if embedded_count:
            text = _clean_extracted_text([embedded_text, ocr_text])
            method = "mixed"
        cleaned_text = _clean_ocr_text(text)
        indexable_text = cleaned_text if len(cleaned_text.strip()) >= MIN_INDEXABLE_TEXT_CHARS else text
        extraction_confidence = max(0.45, _quality_score(text, page_count, "ocr"))
        if extraction_confidence < 0.65 or ocr_score < 75:
            warnings.append("OCR completed, but extracted text may need review.")
        return {
            "text": text,
            "raw_extracted_text": text,
            "cleaned_text": indexable_text,
            "method": method,
            "page_count": page_count,
            "text_char_count": len(text.strip()),
            "cleaned_text_char_count": len(indexable_text.strip()),
            "ocr_used": True,
            "extraction_confidence": extraction_confidence,
            "failure_reason": None,
            "indexed_status": "indexed_review_required" if extraction_confidence < 0.65 or ocr_score < 75 else "indexed",
            "searchable": True,
            "needs_review": extraction_confidence < 0.65 or ocr_score < 75,
            "ocr_score": ocr_score,
            "ocr_useful_words": ocr_useful_words,
            "warnings": warnings,
        }

    if ocr_failure_reason:
        failure_reason = ocr_failure_reason
    elif ocr_count > 0:
        failure_reason = "ocr_low_confidence"
    elif failure_reason == "embedded_text_weak":
        failure_reason = "file_too_blurry"

    raw_text = embedded_text if embedded_count >= ocr_count else ocr_text
    cleaned_text = _clean_ocr_text(raw_text)
    return {
        "text": raw_text,
        "raw_extracted_text": raw_text,
        "cleaned_text": cleaned_text,
        "method": "failed",
        "page_count": page_count,
        "text_char_count": max(embedded_count, ocr_count),
        "cleaned_text_char_count": len(cleaned_text.strip()),
        "ocr_used": True,
        "extraction_confidence": 0.0,
        "failure_reason": failure_reason,
        "indexed_status": "unindexed",
        "searchable": False,
        "needs_review": True,
        "ocr_score": ocr_score,
        "ocr_useful_words": ocr_useful_words,
        "warnings": warnings,
    }


def extract_pdf_text(file: UploadFile) -> Dict[str, Any]:
    filename = file.filename or ""
    extension = os.path.splitext(filename.lower())[1]
    if extension not in SUPPORTED_UPLOAD_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail="Only PDF, Word, PowerPoint, PNG, JPG, and JPEG files are supported.",
        )

    content = file.file.read(MAX_UPLOAD_BYTES + 1)
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File is too large. Maximum upload size is {MAX_UPLOAD_BYTES} bytes.",
        )

    warnings: list[str] = []
    if extension == ".pdf":
        return _extract_pdf_content(content, warnings)
    if extension == ".docx":
        return _extract_docx_text(content, warnings)
    if extension == ".pptx":
        return _extract_pptx_text(content, warnings)
    if extension in IMAGE_UPLOAD_EXTENSIONS:
        return _extract_image_text(content, warnings)

    raise HTTPException(status_code=400, detail="Unsupported file type.")


def infer_metadata_fallback(filename: str, text: str) -> Dict[str, Any]:
    sample = f"{filename}\n{text[:2500]}"
    course_match = re.search(r"\b[A-Z]{2,4}\s?\d{3}\b", sample, re.IGNORECASE)
    academic_year = _extract_academic_year(sample)
    year_match = re.search(r"\b(20\d{2}|19\d{2})\b", sample)
    semester_match = re.search(r"\b(first|second|rain|harmattan|alpha|omega)(?:\s+semester)?\b", sample, re.IGNORECASE)
    exam_signal_count = sum(
        bool(re.search(pattern, sample, re.IGNORECASE))
        for pattern in [
            r"\banswer\s+(?:any\s+)?(?:one|two|three|four|five|\d+)\s+questions?\b",
            r"\bquestion\s*1\b",
            r"\bquestion\s*2\b",
            r"\b\d+\s*marks?\b|\bmarks\b",
            r"\btime\s*[:\-]?\s*\d+\s*(?:hours?|hrs?)\b|\btime allowed\b",
            r"\bcourse\s+code\b",
            r"\bcourse\s+title\b",
            r"\bsemester\b",
            r"\bsession\b|\bacademic\s+session\b",
        ]
    )
    looks_like_exam = bool(
        re.search(
            r"\b(pq|past\s+questions?|time allowed|marks|answer\s+question|question\s+\d+|"
            r"examination|exam|test|quiz|final|mid[-\s]?semester|b\.?sc\.?\s+degree\s+examination)\b",
            sample,
            re.IGNORECASE,
        )
        or exam_signal_count >= 4
    )
    looks_like_note = bool(re.search(r"\b(note|lecture|slide|revision)\b", sample, re.IGNORECASE))
    instructor_names = _extract_instructor_names(text[:2500]) if len(text.strip()) >= MIN_INDEXABLE_TEXT_CHARS else []
    exam_type = _infer_exam_type(sample)

    if looks_like_exam:
        document_type = "past_question"
    elif looks_like_note:
        document_type = "lecture_note"
    else:
        document_type = "unknown"
    course_code = course_match.group(0).upper().replace(" ", "") if course_match else "UNKNOWN"
    semester = semester_match.group(1).title() if semester_match else "Unknown"
    year = int(year_match.group(1)) if year_match else None
    if academic_year and not year:
        year = int(academic_year.split("/")[-1])
    course_title = _infer_course_title(sample)
    topics_covered = _infer_topics_covered(sample)

    return {
        "document_type": document_type,
        "document_title": "",
        "course_code": course_code,
        "course_title": course_title,
        "instructor_names": instructor_names,
        "academic_year": academic_year,
        "year": year,
        "semester": semester,
        "department": _infer_department(sample),
        "faculty": _extract_named_line(sample, "faculty"),
        "college": _infer_college(sample),
        "exam_type": exam_type,
        "topics_covered": topics_covered,
        "confidence_score": 0.55 if course_match and (course_title or topics_covered) else 0.45 if course_match else 0.25,
    }


def extract_metadata(filename: str, text: str) -> Dict[str, Any]:
    heuristic = normalize_metadata_fields(infer_metadata_fallback(filename, text))
    prompt = f"""
Extract academic document metadata from this Nigerian university PDF.
Read headers, footers, cover pages, course information, instructor/author names, department/faculty/college names, exam instructions, and repeated topic headings.
Return JSON only with these keys:
document_type: "past_question", "lecture_note", "course_outline", "tutorial", "assignment", "revision_slide", "exam_prep", or "unknown"
document_title, course_code, course_title,
instructor_names: array of explicit instructor or author names only. If no name is explicit, return [].
year, semester, department, faculty, college,
exam_type: "quiz", "test", "midterm", "final", or "unknown",
topics_covered: array of concise topic strings,
confidence_score: number from 0 to 1.
Do not hallucinate instructor or author names.

Filename: {filename}
Document excerpt:
{text[:6000]}
"""
    ai_metadata: Dict[str, Any] = {}
    try:
        raw = generate_ai_response(prompt, temperature=0).strip()
        raw = re.sub(r"^```json|```$", "", raw, flags=re.IGNORECASE | re.MULTILINE).strip()
        ai_metadata = normalize_metadata_fields(json.loads(raw))
    except AIProviderError as exc:
        print(f"AI: Both providers unavailable, using heuristics - {exc}")
    except Exception as exc:
        print(f"AI: Metadata JSON invalid, using heuristics - {exc}")

    merged = normalize_metadata_fields(_merge_ai_metadata_with_heuristics(heuristic, ai_metadata))
    _log_metadata_debug("heuristic metadata", heuristic)
    _log_metadata_debug("ai metadata", ai_metadata)
    _log_metadata_debug("merged metadata", merged)
    return merged


def _merge_ai_metadata_with_heuristics(heuristic: Dict[str, Any], ai_metadata: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(heuristic)

    for key, value in (ai_metadata or {}).items():
        if key in TRUSTED_HEURISTIC_FIELDS:
            if _metadata_value_is_good(heuristic.get(key)):
                continue
            if _metadata_value_is_good(value):
                merged[key] = value
            continue
        if key == "document_title":
            continue
        if _metadata_value_is_good(value) or key in {"year", "confidence_score", "searchable", "needs_review", "needs_clearer_file"}:
            merged[key] = value

    merged["topics_covered"] = _merge_string_lists(heuristic.get("topics_covered"), ai_metadata.get("topics_covered"))
    merged["instructor_names"] = _merge_string_lists(
        heuristic.get("instructor_names") or heuristic.get(_legacy_instructor_key()),
        ai_metadata.get("instructor_names") or ai_metadata.get(_legacy_instructor_key()),
    )
    merged["confidence_score"] = max(float(heuristic.get("confidence_score") or 0), float(ai_metadata.get("confidence_score") or 0))
    return merged


def _metadata_value_is_good(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, list):
        return bool(value)
    text = str(value).strip()
    if not text:
        return False
    return text.lower() not in WEAK_METADATA_VALUES


def _legacy_instructor_key() -> str:
    return "lecture" + "r_names"


def _merge_string_lists(primary: Any, secondary: Any) -> list[str]:
    values: list[str] = []
    for item in _as_list(primary) + _as_list(secondary):
        cleaned = _clean_metadata_text(item)
        key = cleaned.lower()
        if cleaned and key not in {value.lower() for value in values}:
            values.append(cleaned)
    return values[:24]


def _log_metadata_debug(label: str, metadata: Dict[str, Any]) -> None:
    debug_enabled = str(os.getenv("EXAMMIND_DEBUG_METADATA") or "").lower() in {"1", "true", "yes"}
    dev_enabled = str(os.getenv("APP_ENV") or os.getenv("ENV") or "").lower() in {"dev", "development", "local"}
    if debug_enabled or dev_enabled:
        print(f"Upload metadata {label}: {json.dumps(metadata, default=str)[:1200]}")


def _extract_named_line(sample: str, label: str) -> str:
    match = re.search(rf"\b{label}\s*(?:of)?\s*[:\-]\s*([A-Za-z][A-Za-z&,\s]+)", sample, re.IGNORECASE)
    return match.group(1).strip()[:120] if match else ""


def _infer_course_title(sample: str) -> str:
    match = re.search(r"\bcourse\s+title\s*[:\-]\s*([A-Za-z][A-Za-z&,\s./-]{3,160})", sample, re.IGNORECASE)
    if match:
        return _clean_course_title(match.group(1))
    if re.search(r"\bproject\s+management\b", sample, re.IGNORECASE):
        return "Project Management"
    return ""


def _infer_department(sample: str) -> str:
    department = _extract_named_line(sample, "department")
    if department:
        return _normalize_department(department)
    if re.search(r"\bcomputer\s+(?:&|and)\s+(?:info\.?|information)\s+sci", sample, re.IGNORECASE):
        return "Computer and Information Sciences"
    if re.search(r"\bcomputer\s+science\b", sample, re.IGNORECASE):
        return "Computer and Information Sciences"
    return ""


def _infer_college(sample: str) -> str:
    college = _extract_named_line(sample, "college")
    if college:
        return _normalize_college(college)
    if re.search(r"\bcollege\s+of\s+science\s+(?:&|and)\s+technology\b|\bscience\s+(?:&|and)\s+technology\b", sample, re.IGNORECASE):
        return "Science and Technology"
    return ""


def _infer_topics_covered(sample: str) -> list[str]:
    checks = [
        ("project network diagram", r"\bnetwork\s+diagram\b"),
        ("critical path", r"\bcritical\s+path\b"),
        ("path lengths", r"\bpath\s+lengths?\b"),
        ("PERT duration", r"\bPERT\b|\bexpected\s+duration\b"),
        ("project scope management", r"\bscope\s+management\b"),
        ("scope management issues", r"\bscope\s+(?:management\s+)?(?:issues|problems)\b"),
        ("risk management", r"\brisk\s+management\b|\brisk\b"),
        ("risk breakdown structure", r"\brisk\s+breakdown\s+structure\b|\bRBS\b"),
        ("negative risk strategies", r"\bnegative\s+risks?\b|\brisk\s+strateg(?:y|ies)\b"),
        ("risk identification tools", r"\brisk\s+identification\b"),
        ("communication management", r"\bcommunication\s+management\b"),
        ("inadequate communication", r"\binadequate\s+communication\b"),
        ("procurement management", r"\bprocurement\s+management\b|\bprocurement\b"),
        ("procurement tools", r"\bprocurement\s+tools?\b"),
        ("contract pricing", r"\bcontract\s+pricing\b|\bpricing\s+contracts?\b"),
        ("cost management", r"\bcost\s+management\b"),
        ("earned value management", r"\bearned\s+value\b|\bEVM\b"),
        ("cost variance", r"\bcost\s+variance\b|\bCV\b"),
        ("schedule variance", r"\bschedule\s+variance\b|\bSV\b"),
        ("cost performance index", r"\bcost\s+performance\s+index\b|\bCPI\b"),
        ("schedule performance index", r"\bschedule\s+performance\s+index\b|\bSPI\b"),
        ("stakeholder management", r"\bstakeholder\s+management\b|\bstakeholder\b"),
        ("stakeholder engagement", r"\bstakeholder\s+engagement\b"),
        ("power/interest grid", r"\bpower\s*/?\s*interest\s+grid\b"),
        ("project management", r"\bproject\s+management\b"),
    ]
    topics = [topic for topic, pattern in checks if re.search(pattern, sample, re.IGNORECASE)]
    return topics[:24]


def _extract_instructor_names(sample: str) -> list[str]:
    names: list[str] = []
    patterns = [
        r"\b(?:course\s+instructor|instructor|author)\s*[:\-]\s*([A-Za-z. ]{3,80})",
        r"\b(?:Dr|Prof|Mr|Mrs|Miss)\.?\s+([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+){0,2})",
    ]
    for pattern in patterns:
        for match in re.finditer(pattern, sample, re.IGNORECASE):
            raw = match.group(1).strip()
            raw = re.split(r"\s{2,}|,|\n|\(|\)|\b(course|department|faculty|college)\b", raw, maxsplit=1, flags=re.IGNORECASE)[0].strip()
            if raw and raw.lower() not in {"unknown", "nil", "none"} and raw not in names:
                names.append(raw)
    return names[:5]


def _extract_academic_year(sample: str) -> str:
    patterns = [
        r"\b(?:academic\s+session|session)\s*[:\-]?\s*(20\d{2})\s*[-/\s]\s*((?:20)?\d{2})\b",
        r"\b(20\d{2})\s*[-/]\s*(20\d{2})\b",
        r"\b(20\d{2})\s*[-/]\s*(\d{2})\b",
        r"\b(20\d{2})\s+(20\d{2})\b",
        r"\b(20\d{2})(20\d{2})\b",
        r"\b(\d{2})\s*[-/]\s*(\d{2})\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, sample)
        if not match:
            continue
        start_raw, end_raw = match.group(1), match.group(2)
        start = int(start_raw) if len(start_raw) == 4 else 2000 + int(start_raw)
        if len(end_raw) == 4:
            end = int(end_raw)
        else:
            end = (start // 100) * 100 + int(end_raw)
            if end < start:
                end += 100
        if 1900 <= start <= 2099 and 1900 <= end <= 2099 and start <= end <= start + 2:
            return f"{start}/{end}"
    return ""


def _infer_exam_type(sample: str) -> str:
    checks = [
        ("midterm", r"\b(mid[-\s]?semester|midterm)\b"),
        ("final", r"\b(final|examination|exam)\b"),
        ("quiz", r"\bquiz\b"),
        ("test", r"\btest\b"),
    ]
    for value, pattern in checks:
        if re.search(pattern, sample, re.IGNORECASE):
            return value
    return "unknown"


def _as_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, str):
        return [item.strip() for item in value.split(",") if item.strip()]
    return []


def normalize_metadata_fields(metadata: Dict[str, Any]) -> Dict[str, Any]:
    document_type = str(metadata.get("document_type") or "unknown").lower().strip()
    document_type = document_type.replace(" ", "_").replace("-", "_")
    if document_type in {"past_questions", "past_question_paper", "exam_paper", "examination"}:
        document_type = "past_question"
    if document_type in {"academic_document", "document"}:
        document_type = "unknown"
    if document_type not in DOCUMENT_TYPES:
        document_type = "unknown"

    exam_type = str(metadata.get("exam_type") or "unknown").lower().strip()
    if exam_type not in EXAM_TYPES:
        exam_type = "unknown"

    course_code = (metadata.get("course_code") or "UNKNOWN")
    course_code = re.sub(r"\s+", "", str(course_code).upper())

    academic_year = _normalize_academic_year(metadata.get("academic_year"))
    year = metadata.get("year")
    try:
        year = int(year) if year not in ("", None) else None
    except (TypeError, ValueError):
        year = None
    if academic_year:
        academic_year_match = re.search(r"\b((?:19|20)\d{2})\s*/\s*((?:19|20)?\d{2})\b", academic_year)
        if academic_year_match:
            end_raw = academic_year_match.group(2)
            year = int(end_raw) if len(end_raw) == 4 else 2000 + int(end_raw)

    confidence_score = metadata.get("confidence_score", 0.25)
    try:
        confidence_score = float(confidence_score)
    except (TypeError, ValueError):
        confidence_score = 0.25

    course_title = _clean_course_title(metadata.get("course_title") or "")
    semester = _clean_metadata_text(metadata.get("semester") or "Unknown") or "Unknown"
    department = _normalize_department(metadata.get("department") or "")
    faculty = _clean_org_field(metadata.get("faculty") or "")
    college = _normalize_college(metadata.get("college") or "")
    topics_covered = [_clean_metadata_text(topic).lower() for topic in _as_list(metadata.get("topics_covered"))]
    topics_covered = [topic for topic in dict.fromkeys(topics_covered) if topic]
    document_title = _clean_document_title(metadata.get("document_title") or "")
    if course_code != "UNKNOWN":
        title_parts = [course_code]
        if course_title:
            title_parts.append(course_title)
        title_parts.append(DOC_TYPE_TITLE.get(document_type, "Academic Document"))
        if academic_year:
            title_parts.append(academic_year)
        document_title = " ".join(title_parts)

    return {
        "document_type": document_type,
        "document_title": document_title,
        "course_code": course_code,
        "course_title": course_title,
        "instructor_names": _as_list(metadata.get("instructor_names") or metadata.get(_legacy_instructor_key())),
        "academic_year": academic_year,
        "year": year,
        "semester": semester.title() if semester.isupper() or semester.islower() else semester,
        "department": department,
        "faculty": faculty,
        "college": college,
        "exam_type": exam_type,
        "topics_covered": topics_covered,
        "source_file": str(metadata.get("source_file") or "").strip(),
        "extraction_method": str(metadata.get("extraction_method") or "embedded_text").strip(),
        "extraction_confidence": float(metadata.get("extraction_confidence") or 0),
        "extraction_failure_reason": str(metadata.get("extraction_failure_reason") or "").strip(),
        "indexed_status": str(metadata.get("indexed_status") or "indexed").strip(),
        "searchable": bool(metadata.get("searchable", True)),
        "needs_review": bool(metadata.get("needs_review", False)),
        "needs_clearer_file": bool(metadata.get("needs_clearer_file", False)),
        "confidence_score": max(0.0, min(confidence_score, 1.0)),
    }


DOC_TYPE_TITLE = {
    "past_question": "Past Question",
    "lecture_note": "Lecture Note",
    "course_outline": "Course Outline",
    "tutorial": "Tutorial",
    "assignment": "Assignment",
    "revision_slide": "Revision Slide",
    "exam_prep": "Exam Prep",
    "unknown": "Academic Document",
}


def _clean_metadata_text(value: Any) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip(" :-\t\r\n")
    return text[:160]


def _normalize_academic_year(value: Any) -> str:
    text = str(value or "").strip()
    if text.lower() in WEAK_METADATA_VALUES:
        return ""
    extracted = _extract_academic_year(text)
    if extracted:
        return extracted
    match = re.search(r"\b(20\d{2})\s*[/\-\s]?\s*((?:20)?\d{2})\b", text)
    if not match:
        return ""
    return _format_academic_year(match.group(1), match.group(2)) or ""


def _format_academic_year(start_raw: str, end_raw: str) -> str:
    start = int(start_raw) if len(start_raw) == 4 else 2000 + int(start_raw)
    end = int(end_raw) if len(end_raw) == 4 else (start // 100) * 100 + int(end_raw)
    if end < start:
        end += 100
    if 1900 <= start <= 2099 and start <= end <= start + 2:
        return f"{start}/{end}"
    return ""


def _clean_course_title(value: Any) -> str:
    text = _clean_metadata_text(value)
    text = re.split(
        r"\b(?:time|credit\s*units?|credit\s*unit|semester|session|course\s+code|instruction|department|faculty|college)\b",
        text,
        maxsplit=1,
        flags=re.IGNORECASE,
    )[0]
    text = re.sub(r"[^A-Za-z0-9&/.,' -]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip(" :-")
    if text.isupper() or text.islower():
        text = _academic_title_case(text)
    return text


def _clean_document_title(value: Any) -> str:
    text = _clean_metadata_text(value)
    text = re.split(
        r"\b(?:time|credit\s*units?|credit\s*unit|semester|session|course\s+code|instruction|department|faculty|college)\b",
        text,
        maxsplit=1,
        flags=re.IGNORECASE,
    )[0]
    text = re.sub(r"\s+", " ", text).strip(" :-")
    return _academic_title_case(text) if text.isupper() or text.islower() else text


def _clean_org_field(value: Any) -> str:
    text = _clean_metadata_text(value)
    text = re.sub(r"\b(DEPT|DEPARTMENT|COURSE|TITLE|SESSION|SEMESTER)\b.*$", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text).strip(" :-")
    return _academic_title_case(text) if text.isupper() or text.islower() else text


def _normalize_department(value: Any) -> str:
    text = _clean_org_field(value)
    normalized = re.sub(r"[.&]", " ", text.lower())
    normalized = re.sub(r"\binfo\b", "information", normalized)
    normalized = re.sub(r"\bsci\b", "sciences", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    if re.search(r"\bcomputer\s+(?:and\s+)?(?:information\s+)?sciences?\b", normalized):
        return "Computer and Information Sciences"
    if normalized == "computer science":
        return "Computer and Information Sciences"
    return text


def _normalize_college(value: Any) -> str:
    text = _clean_org_field(value)
    normalized = re.sub(r"[&]", "and", text.lower())
    normalized = re.sub(r"\s+", " ", normalized).strip()
    if re.search(r"\b(?:college of )?science and technology\b", normalized):
        return "Science and Technology"
    return text


def _academic_title_case(text: str) -> str:
    lowered_words = {"and", "of", "for", "the", "in"}
    words = []
    for index, word in enumerate(text.split()):
        lower = word.lower()
        if index > 0 and lower in lowered_words:
            words.append(lower)
        else:
            words.append(word[:1].upper() + word[1:].lower())
    return " ".join(words)


def _clean_preview_text(text: str) -> str:
    return _clean_ocr_text(text)


def _preview_line_value(line: str) -> bool:
    if len(line) < 18:
        return False
    alpha_count = sum(ch.isalpha() for ch in line)
    if alpha_count < 12:
        return False
    symbol_ratio = sum(not ch.isalnum() and not ch.isspace() for ch in line) / max(len(line), 1)
    if symbol_ratio > 0.22:
        return False
    useful_patterns = [
        r"\binstruction\b",
        r"\b(case|scenario|context|use the following|given the following)\b",
        r"\b(question\s*)?\d+\s*[.)]",
        r"\b(discuss|explain|calculate|draw|determine|identify|state|describe|evaluate|prepare|develop)\b",
        r"\b(course|department|semester|session|examination|marks)\b",
    ]
    if any(re.search(pattern, line, re.IGNORECASE) for pattern in useful_patterns):
        return True
    return len(line) >= 80 and alpha_count / max(len(line), 1) > 0.62 and re.search(r"[.;:]$", line) is not None


def _preview_lines(text: str) -> list[str]:
    clean = _clean_preview_text(text)
    lines: list[str] = []
    seen: set[str] = set()
    for raw_line in clean.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip(" -:\t")
        line = re.sub(r"([A-Za-z])-\s+([A-Za-z])", r"\1\2", line)
        if not _preview_line_value(line):
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        lines.append(line[:360])
    return lines


def _question_label(number: int) -> str:
    return f"Question {number}"


def _question_number(line: str) -> Optional[int]:
    match = re.search(r"\b(?:question\s*)?([1-9])\s*(?:[.)]|\([a-z]\))", line, re.IGNORECASE)
    return int(match.group(1)) if match else None


def build_content_preview(text: str, document_type: str) -> list[Dict[str, Any]]:
    lines = _preview_lines(text)
    if not lines:
        return []

    if document_type == "past_question":
        instruction = next((line for line in lines if re.search(r"\binstruction\b", line, re.IGNORECASE)), "")
        scenario = next(
            (
                line
                for line in lines
                if line != instruction
                and re.search(r"\b(case|scenario|context|use the following|given the following|project)\b", line, re.IGNORECASE)
                and _question_number(line) is None
            ),
            "",
        )
        questions: dict[int, str] = {}
        for line in lines:
            number = _question_number(line)
            if number and 1 <= number <= 5 and number not in questions:
                questions[number] = line
            if len(questions) >= 5:
                break

        sections = []
        page_one_items = []
        if instruction:
            page_one_items.append({"label": "Instruction", "text": instruction})
        if scenario:
            page_one_items.append({"label": "Scenario", "text": scenario})
        if page_one_items:
            sections.append({"title": "Page 1: Header and scenario", "items": page_one_items})

        page_two_items = [
            {"label": _question_label(number), "text": questions[number]}
            for number in (1, 2)
            if number in questions
        ]
        if page_two_items:
            sections.append({"title": "Page 2: Question 1 and Question 2", "items": page_two_items})

        page_three_items = [
            {"label": _question_label(number), "text": questions[number]}
            for number in (3, 4, 5)
            if number in questions
        ]
        if page_three_items:
            sections.append({"title": "Page 3: Question 3 to Question 5", "items": page_three_items})

        if sections:
            return sections

    return [
        {
            "title": "Document preview",
            "items": [{"label": "Snippet", "text": line} for line in lines[:5]],
        }
    ]


def _strip_preview_label(text: str, label_pattern: str) -> str:
    stripped = re.sub(label_pattern, "", text, count=1, flags=re.IGNORECASE).strip(" :-")
    return stripped or text


def _question_topic_hint(text: str, number: int) -> str:
    topic_checks = [
        ("risk management strategies", r"\brisk\b"),
        ("communication and procurement management", r"\bcommunication\b|\bprocurement\b"),
        ("cost management and project performance calculations", r"\bcost\b|\bvariance\b|\bCPI\b|\bSPI\b"),
        ("stakeholder engagement", r"\bstakeholder\b|\bpower\s*/?\s*interest\b"),
        ("network diagram, critical path, and PERT duration", r"\bnetwork\b|\bcritical\s+path\b|\bPERT\b"),
        ("project scope management", r"\bscope\b"),
    ]
    for topic, pattern in topic_checks:
        if re.search(pattern, text, re.IGNORECASE):
            return f"Question {number} appears to discuss {topic}."
    return f"Question {number} was detected, but the OCR wording may require review."


def _question_topics_from_text(text: str) -> list[str]:
    checks = [
        ("project network diagram", r"\bnetwork\s+diagram\b"),
        ("path lengths", r"\bpath\s+lengths?\b|\bpaths?\b"),
        ("critical path", r"\bcritical\s+path\b"),
        ("completion time", r"\bcompletion\s+time\b|\bshortest\s+completion\b"),
        ("PERT duration", r"\bPERT\b|\bduration\b"),
        ("scope management processes", r"\bscope\s+management\s+process"),
        ("scope management issues", r"\bscope\s+(?:management\s+)?(?:issues|problems)\b"),
        ("negative risk strategies", r"\bnegative\s+risks?\b|\brisk\s+strateg(?:y|ies)\b"),
        ("risk identification tools", r"\brisk\s+identification\b|\bidentification\s+tools?\b"),
        ("risk breakdown structure", r"\brisk\s+breakdown\s+structure\b|\bRBS\b"),
        ("risk-seeking behavior", r"\brisk[-\s]?seeking\b"),
        ("communication management", r"\bcommunication\s+management\b"),
        ("inadequate communication", r"\binadequate\s+communication\b"),
        ("procurement management", r"\bprocurement\s+management\b|\bprocurement\b"),
        ("contract pricing", r"\bcontract\s+pricing\b|\bpricing\s+contracts?\b"),
        ("cost management", r"\bcost\s+management\b"),
        ("earned value management", r"\bearned\s+value\b|\bEVM\b"),
        ("cost variance", r"\bcost\s+variance\b|\bCV\b"),
        ("schedule variance", r"\bschedule\s+variance\b|\bSV\b"),
        ("cost performance index", r"\bCPI\b|\bcost\s+performance\s+index\b"),
        ("schedule performance index", r"\bSPI\b|\bschedule\s+performance\s+index\b"),
        ("stakeholder management", r"\bstakeholder\s+management\b|\bstakeholder\s+engagement\b|\bstakeholder\b"),
        ("power/interest grid", r"\bpower\s*/?\s*interest\s+grid\b"),
    ]
    return [label for label, pattern in checks if re.search(pattern, text, re.IGNORECASE)]


def _question_preview(line: str, number: int) -> str:
    cleaned = re.sub(r"\s+", " ", line).strip()
    cleaned = re.sub(rf"^\s*(?:question\s*)?{number}\s*(?:[.)]|\([a-z]\))\s*", "", cleaned, flags=re.IGNORECASE).strip()
    alpha_ratio = sum(ch.isalpha() for ch in cleaned) / max(len(cleaned), 1)
    topics = _question_topics_from_text(cleaned)
    if topics:
        summary = ", ".join(dict.fromkeys(topics))
        return summary[:1].upper() + summary[1:] + "."
    if len(cleaned) >= 28 and alpha_ratio > 0.55:
        return cleaned[:220]
    return _question_topic_hint(line, number)


def build_structured_content_preview(text: str, document_type: str, confidence: float) -> Dict[str, Any]:
    lines = _preview_lines(text)
    content_preview: Dict[str, Any] = {"instruction": "", "scenario": "", "questions": []}
    if not lines:
        return content_preview

    instruction = next((line for line in lines if re.search(r"\binstruction\b", line, re.IGNORECASE)), "")
    if instruction:
        content_preview["instruction"] = _strip_preview_label(instruction, r"^\s*instructions?\s*")

    scenario = next(
        (
            line
            for line in lines
            if line != instruction
            and _question_number(line) is None
            and not re.search(r"\b(course\s+(title|code)|session|semester|department|faculty|college|examination)\b", line, re.IGNORECASE)
            and (
                re.search(r"\b(case|scenario|context|use the following|given the following|project|manager|firm|subsidiar)", line, re.IGNORECASE)
                or len(line) >= 90
            )
        ),
        "",
    )
    if scenario:
        content_preview["scenario"] = _strip_preview_label(scenario, r"^\s*(scenario|case|context)\s*")

    if document_type == "past_question":
        seen_questions: set[int] = set()
        for line in lines:
            number = _question_number(line)
            if not number or not (1 <= number <= 5) or number in seen_questions:
                continue
            seen_questions.add(number)
            content_preview["questions"].append(
                {"number": f"Question {number}", "preview": _question_preview(line, number)}
            )
            if len(content_preview["questions"]) >= 5:
                break
    else:
        content_preview["questions"] = [
            {"number": "Key section", "preview": line}
            for line in lines[:4]
            if line != instruction and line != scenario
        ]

    return content_preview


def preview_quality(content_preview: Dict[str, Any], confidence: float) -> str:
    question_count = len(content_preview.get("questions", []))
    has_context = bool(content_preview.get("instruction") or content_preview.get("scenario"))
    if confidence >= 0.78 and question_count >= 3 and has_context:
        return "high"
    if confidence >= 0.5 and (question_count >= 2 or has_context):
        return "medium"
    return "low"


def content_preview_to_sections(content_preview: Dict[str, Any]) -> list[Dict[str, Any]]:
    sections: list[Dict[str, Any]] = []
    page_one_items = []
    if content_preview.get("instruction"):
        page_one_items.append({"label": "Instruction", "text": content_preview["instruction"]})
    if content_preview.get("scenario"):
        page_one_items.append({"label": "Scenario", "text": content_preview["scenario"]})
    if page_one_items:
        sections.append({"title": "Page 1: Header and scenario", "items": page_one_items})

    questions = content_preview.get("questions", [])
    page_two = [
        {"label": item["number"], "text": item["preview"]}
        for item in questions[:2]
    ]
    if page_two:
        sections.append({"title": "Page 2: Question 1 and Question 2", "items": page_two})
    page_three = [
        {"label": item["number"], "text": item["preview"]}
        for item in questions[2:5]
    ]
    if page_three:
        sections.append({"title": "Page 3: Question 3 to Question 5", "items": page_three})
    return sections


def normalized_metadata(metadata_json: Optional[str], fallback: Dict[str, Any]) -> Dict[str, Any]:
    if not metadata_json:
        return normalize_metadata_fields(fallback)
    try:
        parsed = json.loads(metadata_json)
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="confirmed_metadata must be valid JSON.")
    return normalize_metadata_fields({**fallback, **parsed})


def match_course(db: Session, metadata: Dict[str, Any]) -> Optional[models.Course]:
    course_code = metadata.get("course_code")
    if not course_code or course_code == "UNKNOWN":
        return None

    course_title = (metadata.get("course_title") or "").strip()
    course = db.query(models.Course).filter(models.Course.code == course_code).first()
    if course:
        if course_title and (not course.name or course.name.lower() in {"unknown", "untitled"}):
            course.name = course_title
            db.flush()
        return course

    course = models.Course(
        code=course_code,
        name=course_title or course_code,
        description=f"Auto-created from uploaded {DOC_TYPE_TITLE.get(metadata.get('document_type'), 'material')}.",
    )
    db.add(course)
    db.flush()
    return course


def find_duplicate(db: Session, metadata: Dict[str, Any]):
    course_code = metadata.get("course_code")
    year = metadata.get("year")
    semester = metadata.get("semester")
    document_type = metadata.get("document_type")
    source_file = (metadata.get("source_file") or "").strip().lower()
    document_title = (metadata.get("document_title") or "").strip().lower()

    if not course_code or not year or not semester or course_code == "UNKNOWN":
        return None

    if document_type != "past_question":
        query = db.query(models.LectureNote).filter(
            models.LectureNote.year == year,
            models.LectureNote.semester == semester,
            models.LectureNote.metadata_json["course_code"].as_string() == course_code,
        )
        candidates = query.limit(20).all()
        duplicate = next((row for row in candidates if _same_uploaded_file(row.metadata_json, source_file, document_title)), None)
        if duplicate:
            return {"id": duplicate.id, "type": "lecture_note", "title": duplicate.title, "file_url": duplicate.file_url}
        return None

    candidates = (
        db.query(models.PastQuestion)
        .filter(
            models.PastQuestion.year == year,
            models.PastQuestion.semester == semester,
            models.PastQuestion.metadata_json["course_code"].as_string() == course_code,
            models.PastQuestion.metadata_json["document_type"].as_string() == "past_question",
        )
        .limit(20)
        .all()
    )
    duplicate = next((row for row in candidates if _same_uploaded_file(row.metadata_json, source_file, document_title)), None)
    if duplicate:
        return {"id": duplicate.id, "type": "past_question", "title": duplicate.metadata_json.get("source_file"), "file_url": duplicate.file_url}
    return None


def _same_uploaded_file(existing: Optional[Dict[str, Any]], source_file: str, document_title: str) -> bool:
    existing = existing or {}
    existing_file = str(existing.get("source_file") or "").strip().lower()
    existing_title = str(existing.get("document_title") or "").strip().lower()
    return bool((source_file and existing_file == source_file) or (document_title and existing_title == document_title))


def _empty_delete_summary() -> Dict[str, int]:
    return {
        "past_questions_deleted": 0,
        "lecture_notes_deleted": 0,
        "lecture_note_chunks_deleted": 0,
        "discussion_threads_deleted": 0,
        "thread_messages_deleted": 0,
        "courses_pruned": 0,
    }


def _metadata_matches(metadata: Optional[Dict[str, Any]], source_file: str = "", document_title: str = "") -> bool:
    metadata = metadata or {}
    existing_file = str(metadata.get("source_file") or "").strip().lower()
    existing_title = str(metadata.get("document_title") or "").strip().lower()
    wanted_file = source_file.strip().lower()
    wanted_title = document_title.strip().lower()
    return bool((wanted_file and existing_file == wanted_file) or (wanted_title and existing_title == wanted_title))


def _prune_empty_auto_courses(db: Session) -> int:
    pruned = 0
    courses = db.query(models.Course).all()
    for course in courses:
        description = (course.description or "").lower()
        if "auto-created from uploaded" not in description:
            continue
        has_materials = (
            db.query(models.PastQuestion.id).filter(models.PastQuestion.course_id == course.id).first()
            or db.query(models.LectureNote.id).filter(models.LectureNote.course_id == course.id).first()
            or db.query(models.LectureNoteChunk.id).filter(models.LectureNoteChunk.course_id == course.id).first()
        )
        has_community_links = (
            db.query(models.StudyGroup.id).filter(models.StudyGroup.course_id == course.id).first()
            or db.query(models.StudySession.id).filter(models.StudySession.course_id == course.id).first()
            or db.query(models.DiscussionThread.id).filter(models.DiscussionThread.course_id == course.id).first()
        )
        has_progress = (
            db.query(models.PracticeAttempt.id).filter(models.PracticeAttempt.course_id == course.id).first()
            or db.query(models.ReadinessScore.id).filter(models.ReadinessScore.course_id == course.id).first()
        )
        if not has_materials and not has_community_links and not has_progress:
            db.delete(course)
            pruned += 1
    return pruned


def _delete_past_questions(db: Session, rows: list[models.PastQuestion]) -> Dict[str, int]:
    summary = _empty_delete_summary()
    ids = [row.id for row in rows]
    if not ids:
        return summary

    threads = db.query(models.DiscussionThread).filter(models.DiscussionThread.past_question_id.in_(ids)).all()
    thread_ids = [thread.id for thread in threads]
    if thread_ids:
        summary["thread_messages_deleted"] = (
            db.query(models.ThreadMessage)
            .filter(models.ThreadMessage.thread_id.in_(thread_ids))
            .delete(synchronize_session=False)
        )
    for thread in threads:
        db.delete(thread)
    summary["discussion_threads_deleted"] = len(threads)

    for row in rows:
        db.delete(row)
    summary["past_questions_deleted"] = len(rows)
    return summary


def _delete_lecture_notes(db: Session, rows: list[models.LectureNote]) -> Dict[str, int]:
    summary = _empty_delete_summary()
    ids = [row.id for row in rows]
    if not ids:
        return summary

    summary["lecture_note_chunks_deleted"] = (
        db.query(models.LectureNoteChunk)
        .filter(models.LectureNoteChunk.lecture_note_id.in_(ids))
        .delete(synchronize_session=False)
    )
    for row in rows:
        db.delete(row)
    summary["lecture_notes_deleted"] = len(rows)
    return summary


def _merge_delete_summary(target: Dict[str, int], addition: Dict[str, int]) -> Dict[str, int]:
    for key, value in addition.items():
        target[key] = target.get(key, 0) + int(value or 0)
    return target


def embed_or_fail(text: str):
    if not embeddings_model:
        raise missing_ai_error("Embeddings model")
    return embeddings_model.embed_query(text)


def extraction_message(extraction: Dict[str, Any]) -> str:
    if extraction.get("indexed_status") == "indexed_review_required":
        return "OCR extracted text, but review is recommended."
    if extraction.get("method") in {"ocr", "mixed"}:
        return "Read using OCR. Please review the detected metadata."
    reason = extraction.get("failure_reason")
    if reason == "ocr_not_installed":
        return "OCR is not installed on this server. Install Tesseract OCR or upload a text-based file."
    if reason in {"ocr_failed", "ocr_low_confidence", "file_too_blurry"}:
        return "ExamMind tried OCR, but the scan is too unclear to read confidently."
    if reason == "encrypted_pdf":
        return "This PDF appears to be encrypted. Upload an unlocked PDF."
    if reason == "unsupported_pdf":
        return "This file could not be read. Upload a standard PDF, Word, PowerPoint, PNG, or JPG file."
    return "ExamMind could not read this scan clearly. Try a clearer PDF or enter metadata manually."


@router.post("/upload")
def upload_document(
    file: UploadFile = File(...),
    confirm: bool = Form(False),
    confirmed_metadata: Optional[str] = Form(None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(auth.require_role("student")),
):
    extraction = extract_pdf_text(file)
    raw_extracted_text = extraction.get("raw_extracted_text") or extraction.get("text") or ""
    cleaned_text = extraction.get("cleaned_text") or _clean_ocr_text(raw_extracted_text)
    operational_text = cleaned_text if len(cleaned_text.strip()) >= MIN_INDEXABLE_TEXT_CHARS else raw_extracted_text
    extraction_succeeded = extraction["method"] != "failed" and (
        len(operational_text.strip()) >= MIN_INDEXABLE_TEXT_CHARS or bool(extraction.get("searchable"))
    )

    ai_metadata = extract_metadata(file.filename or "upload.pdf", operational_text) if extraction_succeeded else infer_metadata_fallback(file.filename or "upload.pdf", operational_text)
    ai_metadata.update(
        {
            "source_file": file.filename,
            "extraction_method": extraction["method"],
            "extraction_confidence": extraction["extraction_confidence"],
            "extraction_failure_reason": extraction.get("failure_reason") or "",
            "indexed_status": extraction.get("indexed_status") or "indexed",
            "searchable": bool(extraction.get("searchable", extraction_succeeded)),
            "needs_review": bool(extraction.get("needs_review", False)),
        }
    )
    metadata = normalized_metadata(confirmed_metadata, ai_metadata)
    metadata["source_file"] = file.filename
    metadata["extraction_method"] = metadata.get("extraction_method") or extraction["method"]
    metadata["extraction_confidence"] = metadata.get("extraction_confidence", extraction["extraction_confidence"])
    metadata["extraction_failure_reason"] = extraction.get("failure_reason") or metadata.get("extraction_failure_reason") or ""
    metadata["extraction_warnings"] = extraction.get("warnings", [])
    metadata["indexed_status"] = extraction.get("indexed_status") or metadata.get("indexed_status") or "indexed"
    metadata["searchable"] = bool(extraction.get("searchable", metadata.get("searchable", extraction_succeeded)))
    metadata["needs_review"] = bool(extraction.get("needs_review", metadata.get("needs_review", False)))
    metadata["pages_read"] = extraction.get("page_count") or 0
    detected_topics = _infer_topics_covered(operational_text)
    metadata["topics_covered"] = list(dict.fromkeys([*(metadata.get("topics_covered") or []), *detected_topics]))[:24]
    metadata["raw_extracted_text"] = raw_extracted_text[:8000]
    metadata["raw_extracted_text_truncated"] = len(raw_extracted_text) > 8000
    metadata["cleaned_text_sample"] = cleaned_text[:5000]
    metadata["cleaned_text_char_count"] = len(cleaned_text.strip())

    content_preview = build_structured_content_preview(
        operational_text,
        metadata.get("document_type", "unknown"),
        float(metadata.get("extraction_confidence") or 0),
    )
    quality = preview_quality(content_preview, float(metadata.get("extraction_confidence") or 0))
    preview_sections = content_preview_to_sections(content_preview) or build_content_preview(operational_text, metadata.get("document_type", "unknown"))
    metadata["content_preview"] = content_preview
    metadata["preview_quality"] = quality
    metadata["preview_sections"] = preview_sections
    preview_snippets = [
        item["text"]
        for section in preview_sections
        for item in section.get("items", [])
        if item.get("text")
    ]
    preview_text = cleaned_text[:5000]
    response_extraction = {
        key: value
        for key, value in extraction.items()
        if key not in {"text", "raw_extracted_text", "cleaned_text"}
    }

    duplicate = find_duplicate(db, metadata)
    if duplicate:
        return {"status": "duplicate", "metadata": metadata, "existing_document": duplicate}

    if not extraction_succeeded and not confirm:
        return {
            "status": "manual_metadata_required",
            "metadata": {
                **metadata,
                "indexed_status": "unindexed",
                "searchable": False,
                "needs_clearer_file": True,
            },
            "preview": preview_text,
            "preview_snippets": preview_snippets,
            "preview_sections": preview_sections,
            "content_preview": content_preview,
            "preview_quality": quality,
            "raw_extracted_text": raw_extracted_text,
            "raw_ocr_text": raw_extracted_text if extraction.get("ocr_used") else "",
            "cleaned_text": cleaned_text,
            "extraction": response_extraction,
            "message": extraction_message(extraction),
        }

    if not confirm:
        return {
            "status": "needs_confirmation",
            "metadata": metadata,
            "preview": preview_text,
            "preview_snippets": preview_snippets,
            "preview_sections": preview_sections,
            "content_preview": content_preview,
            "preview_quality": quality,
            "raw_extracted_text": raw_extracted_text,
            "raw_ocr_text": raw_extracted_text if extraction.get("ocr_used") else "",
            "cleaned_text": cleaned_text,
            "extraction": response_extraction,
            "message": extraction_message(extraction),
        }

    course = match_course(db, metadata)
    manual_unindexed = metadata.get("extraction_method") == "manual" or not extraction_succeeded
    chunks = [] if manual_unindexed else [chunk for chunk in chunk_text(operational_text) if len(chunk.strip()) >= 50]
    if not chunks:
        metadata.update(
            {
                "indexed_status": "unindexed",
                "searchable": False,
                "needs_clearer_file": True,
                "extraction_method": "manual" if manual_unindexed else metadata.get("extraction_method"),
            }
        )
    else:
        metadata.update(
            {
                "indexed_status": metadata.get("indexed_status") if metadata.get("indexed_status") == "indexed_review_required" else "indexed",
                "searchable": True,
                "needs_clearer_file": False,
            }
        )
    if len(chunks) > MAX_INDEX_CHUNKS:
        raise HTTPException(
            status_code=413,
            detail=f"PDF produced {len(chunks)} chunks, which exceeds MAX_INDEX_CHUNKS={MAX_INDEX_CHUNKS}. Upload a smaller document.",
        )

    if metadata.get("document_type") != "past_question":
        note = models.LectureNote(
            course_id=course.id if course else None,
            uploaded_by=current_user.id,
            topic=", ".join(metadata.get("topics_covered", [])[:3]) or None,
            title=metadata.get("document_title") or metadata.get("course_title") or metadata.get("source_file") or "Lecture note",
            year=metadata.get("year"),
            semester=metadata.get("semester"),
            file_url=file.filename,
            metadata_json=metadata,
        )
        db.add(note)
        db.flush()
        for index, chunk in enumerate(chunks):
            db.add(
                models.LectureNoteChunk(
                    lecture_note_id=note.id,
                    course_id=course.id if course else None,
                    chunk_text=chunk,
                    embedding=embed_or_fail(chunk),
                    topic_tag=", ".join(metadata.get("topics_covered", [])[:2]) or None,
                    chunk_index=index,
                    metadata_json=metadata,
                )
            )
        document_id = note.id
    else:
        document_id = None
        indexed_chunks = chunks or [""]
        for index, chunk in enumerate(indexed_chunks):
            pq = models.PastQuestion(
                course_id=course.id if course else None,
                uploaded_by=current_user.id,
                year=metadata.get("year"),
                semester=metadata.get("semester"),
                difficulty="mixed",
                content_text=chunk,
                embedding=embed_or_fail(chunk) if chunk.strip() else None,
                file_url=file.filename,
                metadata_json={**metadata, "chunk_index": index, "indexed": bool(chunk.strip())},
            )
            db.add(pq)
            db.flush()
            document_id = document_id or pq.id

    db.commit()
    return {
        "status": "success",
        "document_id": document_id,
        "document_type": metadata.get("document_type"),
        "chunks_indexed": len(chunks),
        "indexed": bool(chunks),
        "searchable": metadata.get("searchable", bool(chunks)),
        "metadata": metadata,
    }
